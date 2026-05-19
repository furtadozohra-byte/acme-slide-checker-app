import streamlit as st
from pptx import Presentation

# --- BASIC CONFIG & PASSWORD GATE ---

APP_PASSWORD = "acme-demo"  # password you'll send them in the email

st.set_page_config(page_title="ACME Slide Checker", layout="wide")

st.title("ACME Slide Checker")

password = st.text_input("Enter password to continue", type="password")
if password != APP_PASSWORD:
    st.warning("Enter the password to use this app.")
    st.stop()

st.write("Automated brand‑compliance checks for ACME presentations.")

# --- FILE UPLOAD ---

uploaded = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

if uploaded:
    try:
        prs = Presentation(uploaded)
        st.success("File loaded successfully!")
        st.write(f"Slides detected: {len(prs.slides)}")

        issues = []

        # RULE 1 — Slide must have a title
        for i, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else ""
            if not title.strip():
                issues.append(f"Slide {i+1}: Missing title")

        # RULE 2 — Title must be Title Case
        def is_title_case(text):
            return text == text.title()

        for i, slide in enumerate(prs.slides):
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                if title and not is_title_case(title):
                    issues.append(f"Slide {i+1}: Title not in Title Case")

        # RULE 3 — No more than 6 lines of text (per text frame)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                    if len(lines) > 6:
                        issues.append(f"Slide {i+1}: Too many lines of text in a block ({len(lines)})")
                        break

        # RULE 4 — No more than 40 words total per slide
        for i, slide in enumerate(prs.slides):
            word_count = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        word_count += len(p.text.split())
            if word_count > 40:
                issues.append(f"Slide {i+1}: Too many words on slide ({word_count})")

        # RULE 5 — No red text (safe colour handling)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            color = run.font.color
                            if color is not None and hasattr(color, "rgb") and color.rgb:
                                if str(color.rgb) == "FF0000":
                                    issues.append(f"Slide {i+1}: Contains red text")
                                    break

        # RULE 6 — Logo or image must appear (simple heuristic)
        for i, slide in enumerate(prs.slides):
            has_image = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = PICTURE
            if not has_image:
                issues.append(f"Slide {i+1}: No logo or image detected")

        # RULE 7 — Slide numbers must exist (simple heuristic)
        for i, slide in enumerate(prs.slides):
            has_number = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.isdigit():
                        has_number = True
                        break
            if not has_number:
                issues.append(f"Slide {i+1}: Missing slide number")

        # --- OUTPUT / UX ---

        st.subheader("Results")

        if issues:
            st.error(f"Issues found: {len(issues)}")

            # Group issues by slide label "Slide X"
            slides_dict = {}
            for issue in issues:
                slide_label = issue.split(":")[0]  # "Slide X"
                slides_dict.setdefault(slide_label, []).append(issue)

            for slide_label, slide_issues in slides_dict.items():
                with st.expander(slide_label, expanded=False):
                    for issue in slide_issues:
                        st.checkbox(issue, value=False, key=issue)
        else:
            st.success("No issues found! This deck meets the current ACME rules.")

    except Exception as e:
        st.error(f"Could not read PPTX: {e}")




